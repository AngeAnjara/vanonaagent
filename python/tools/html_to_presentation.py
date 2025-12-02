import json
import os
import time
from typing import Any

from python.helpers.tool import Tool, Response
from python.helpers.print_style import PrintStyle
from python.helpers.errors import RepairableException
from python.helpers import files


class HtmlToPresentation(Tool):
    async def execute(self, **kwargs) -> Response:
        await self.agent.handle_intervention()

        if not getattr(self.agent.config, "ppt_enabled", False):
            raise RepairableException(
                "PowerPoint generation is disabled. Enable in Settings > PowerPoint & PDF Generation."
            )

        html_path: str = (kwargs.get("html_path") or self.args.get("html_path") or "").strip()
        if not html_path:
            raise RepairableException("Missing required argument: html_path")
        if not os.path.isabs(html_path):
            html_path = files.get_abs_path(html_path)
        if not os.path.exists(html_path):
            raise RepairableException("HTML file not found")

        output_format = (kwargs.get("output_format") or self.args.get("output_format") or self.agent.config.ppt_default_format).strip().lower()
        output_name = (kwargs.get("output_name") or self.args.get("output_name") or f"presentation_{int(time.time())}").strip()

        out_dir = files.get_abs_path("outputs", "presentations")
        files.ensure_directory(out_dir)

        if output_format == "pptx":
            try:
                from pptx import Presentation
            except Exception:
                raise RepairableException("Missing library: python-pptx. Install with: pip install python-pptx")

            try:
                from playwright.async_api import async_playwright
            except Exception:
                raise RepairableException("Missing library: playwright. Install with: pip install playwright and run 'playwright install'")

            # Build a list of screenshots: prefer per-slide elements, fallback to full page
            screenshots: list[str] = []
            try:
                # Build file:// URL cross-platform
                from pathlib import Path
                page_url = Path(html_path).resolve().as_uri()
            except Exception:
                page_url = "file://" + html_path.replace("\\", "/")

            try:
                async with async_playwright() as p:
                    browser = await p.chromium.launch()
                    page = await browser.new_page(viewport={"width": 1280, "height": 720})
                    await page.goto(page_url, wait_until="load")
                    # allow charts to render
                    await page.wait_for_timeout(800)
                    # Try to find slide elements
                    slides = await page.query_selector_all(".slide")
                    if slides:
                        for idx, el in enumerate(slides, start=1):
                            s_path = files.get_abs_path("outputs", "presentations", f"{output_name}_slide_{idx}.png")
                            await el.screenshot(path=s_path)
                            screenshots.append(s_path)
                    if not screenshots:
                        # Fallback to full page screenshot
                        s_path = files.get_abs_path("outputs", "presentations", f"{output_name}.png")
                        await page.screenshot(path=s_path, full_page=True)
                        screenshots.append(s_path)
                    await browser.close()
            except Exception as e:
                raise RepairableException(f"Chart rendering failed or Playwright error: {type(e).__name__}: {e}")

            # Assemble PPTX from one or more screenshots
            try:
                prs = Presentation()
                blank_layout = prs.slide_layouts[6]  # blank
                for s_path in screenshots:
                    slide = prs.slides.add_slide(blank_layout)
                    # Fit image to slide; python-pptx uses EMUs, but passing width/height to scale
                    slide_width = prs.slide_width
                    slide_height = prs.slide_height
                    slide.shapes.add_picture(s_path, 0, 0, width=slide_width, height=slide_height)

                pptx_path = files.get_abs_path("outputs", "presentations", f"{output_name}.pptx")
                prs.save(pptx_path)
                msg = json.dumps({"output_path": pptx_path, "format": "pptx"})
                return Response(message=msg, break_loop=True, additional={"output_path": pptx_path})
            except Exception as e:
                raise RepairableException(f"Failed to create PPTX: {type(e).__name__}: {e}")

        elif output_format == "pdf":
            # Prefer weasyprint if available
            pdf_path = files.get_abs_path("outputs", "presentations", f"{output_name}.pdf")
            try:
                from weasyprint import HTML  # type: ignore
                try:
                    HTML(filename=html_path).write_pdf(pdf_path)
                except Exception as e:
                    raise RepairableException(f"Failed to generate PDF with WeasyPrint: {type(e).__name__}: {e}")
            except Exception:
                # fallback to pdfkit if installed
                try:
                    import pdfkit  # type: ignore
                    pdfkit.from_file(html_path, pdf_path)
                except Exception:
                    raise RepairableException("Missing library for PDF: weasyprint or pdfkit. Install with: pip install weasyprint or pip install pdfkit wkhtmltopdf")

            msg = json.dumps({"output_path": pdf_path, "format": "pdf"})
            return Response(message=msg, break_loop=True, additional={"output_path": pdf_path})

        else:
            raise RepairableException("Invalid output_format. Must be 'pptx' or 'pdf'.")
